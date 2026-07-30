import os
import collections
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Define slide size (16:9 Widescreen)
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# Color Scheme (Sleek Dark Theme)
BG_COLOR = RGBColor(15, 23, 42)        # Dark Slate (Slate 900)
CARD_BG_COLOR = RGBColor(30, 41, 59)   # Lighter Slate (Slate 800)
TEXT_WHITE = RGBColor(255, 255, 255)
TEXT_GRAY = RGBColor(148, 163, 184)    # Slate 400
ACCENT_BLUE = RGBColor(14, 165, 233)   # Sky 500
ACCENT_GREEN = RGBColor(34, 197, 94)   # Green 500
ACCENT_ORANGE = RGBColor(249, 115, 22) # Orange 500

def set_slide_background(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def create_title(slide, text):
    title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.5), Inches(11.833), Inches(0.8))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.name = 'Arial'
    p.font.color.rgb = TEXT_WHITE
    return title_box

def add_bullet_points(slide, left, top, width, height, points):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    
    for i, pt in enumerate(points):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = pt["text"]
        p.font.size = Pt(pt.get("size", 16))
        p.font.name = 'Arial'
        p.font.color.rgb = pt.get("color", TEXT_WHITE)
        p.font.bold = pt.get("bold", False)
        
        # Set indent for sub-bullets
        p.level = pt.get("level", 0)
        p.space_after = Pt(10)
    return box

def main():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    
    # Use blank layout
    blank_layout = prs.slide_layouts[6]
    
    # -------------------------------------------------------------------------
    # SLIDE 1: User & Job
    # -------------------------------------------------------------------------
    slide1 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide1, BG_COLOR)
    create_title(slide1, "1. USER & JOB: Ôn tập buổi học VLearn")
    
    points1 = [
        {"text": "Core JTBD (Hành vi & Kỳ vọng chính):", "size": 18, "color": ACCENT_BLUE, "bold": True},
        {"text": "“Khi vừa học xong một buổi nhiều slide, tôi muốn nắm lại toàn bộ nội dung buổi và biết phần nào nằm ở trang nào, để ôn đúng trọng tâm mà không phải lật lại từng trang.”", "size": 16, "color": TEXT_WHITE, "bold": False, "level": 1},
        
        {"text": "Bằng chứng Định lượng Pain Point (Evidence):", "size": 18, "color": ACCENT_BLUE, "bold": True},
        {"text": "• Khảo sát người dùng thật (n=41):", "size": 16, "color": TEXT_WHITE, "bold": True, "level": 1},
        {"text": "- 96.6% học viên gặp khó khăn khi tự ôn tập do tài liệu quá dài, mất nhiều thời gian đọc lại.", "size": 15, "color": TEXT_GRAY, "level": 2},
        {"text": "- 79.3% chê khả năng tóm tắt của Tutor cũ (chỉ tóm tắt trang đơn lẻ, lan man hoặc báo không tìm thấy).", "size": 15, "color": TEXT_GRAY, "level": 2},
        {"text": "- 89.7% sẵn sàng dùng tính năng tóm tắt toàn bộ buổi học ngay lập tức.", "size": 15, "color": TEXT_GRAY, "level": 2},
        
        {"text": "• Khai phá dữ liệu Chatlog thực tế (2,522 dòng):", "size": 16, "color": TEXT_WHITE, "bold": True, "level": 1},
        {"text": "- 29.1% (367/1,261) lượt Tutor trả lời dạng “không tìm thấy” hoặc từ chối do bị neo ở trang hiện tại.", "size": 15, "color": TEXT_GRAY, "level": 2},
        {"text": "- Nhóm câu hỏi tóm tắt/tổng hợp cả buổi học bị bí nhiều nhất: 62.6% (gấp 3 lần câu hỏi trang).", "size": 15, "color": TEXT_GRAY, "level": 2}
    ]
    add_bullet_points(slide1, Inches(0.75), Inches(1.3), Inches(11.833), Inches(5.5), points1)
    
    # -------------------------------------------------------------------------
    # SLIDE 2: Vì sao chọn tính năng này
    # -------------------------------------------------------------------------
    slide2 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide2, BG_COLOR)
    create_title(slide2, "2. LỰA CHỌN TÍNH NĂNG & ĐÁNH GIÁ IMPACT")
    
    points2 = [
        {"text": "Bảng Impact so sánh 3 ứng viên:", "size": 18, "color": ACCENT_BLUE, "bold": True},
        {"text": "1. Ứng viên ①: Câu hỏi tóm tắt buổi (Bị bí 62.6% trong log - 51.7% user chê - Mất >30p/buổi tự lật)", "size": 16, "color": ACCENT_GREEN, "bold": True, "level": 1},
        {"text": "-> QUYẾT ĐỊNH CHỌN: Tác động trực tiếp lên bài toán ôn tập hàng ngày, dễ sửa lỗi bằng cách nới scope RAG.", "size": 15, "color": TEXT_GRAY, "level": 2},
        
        {"text": "2. Ứng viên ②: Câu hỏi vận hành lớp (Deadline, checkpoint, rubric - Bị bí 42.3%)", "size": 16, "color": TEXT_WHITE, "bold": True, "level": 1},
        {"text": "-> TẠM HOÃN: Việc trả lời sai thông tin hạn nộp có thể gây hậu quả trực tiếp (0 điểm). Cần nhiều nguồn sự thật riêng nằm ngoài slide và quy trình kiểm duyệt khắt khe hơn. (Tích hợp một phần vào Cả môn).", "size": 15, "color": TEXT_GRAY, "level": 2},
        
        {"text": "3. Ứng viên ③: Giải thích trang slide đang đọc dễ hiểu hơn", "size": 16, "color": TEXT_WHITE, "bold": True, "level": 1},
        {"text": "-> LOẠI BỎ: Chỉ bị bí 22.9% - Tutor cũ đã làm tương đối ổn định, cải thiện biên nhỏ hơn nhiều so với nhóm ①.", "size": 15, "color": TEXT_GRAY, "level": 2},
        
        {"text": "Lý do loại/chọn bằng con số cụ thể:", "size": 18, "color": ACCENT_BLUE, "bold": True},
        {"text": "Chọn ứng viên ① vì tỷ lệ lỗi chatlog quá lớn (62.6% lỗi). Giải quyết được gốc rễ của pain point bằng một thay đổi duy nhất trong cơ chế quyết định scope RAG (neo trang -> tự nới buổi).", "size": 15, "color": TEXT_WHITE, "level": 1}
    ]
    add_bullet_points(slide2, Inches(0.75), Inches(1.3), Inches(11.833), Inches(5.5), points2)

    # -------------------------------------------------------------------------
    # SLIDE 3: Giải pháp & demo live
    # -------------------------------------------------------------------------
    slide3 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide3, BG_COLOR)
    create_title(slide3, "3. GIẢI PHÁP VÀ KỊCH BẢN DEMO LIVE")
    
    points3 = [
        {"text": "Lát cắt trải nghiệm một câu (Core Slice):", "size": 18, "color": ACCENT_BLUE, "bold": True},
        {"text": "“Học viên đang đọc một trang slide bất kỳ -> bấm hoặc gõ hỏi 'tóm tắt buổi học hôm nay' -> Hệ thống tự quyết định nới phạm vi truy xuất từ trang hiện tại ra toàn bộ slide của buổi -> Trả về cấu trúc 6 mục, takeaway, thuật ngữ kèm link trang bấm được để nhảy trang.”", "size": 16, "color": TEXT_WHITE, "bold": False, "level": 1},
        
        {"text": "Quyết định AI & Model thực hiện:", "size": 18, "color": ACCENT_BLUE, "bold": True},
        {"text": "• AI tự quyết định phạm vi truy xuất (Scope Level) từ câu hỏi -> Trích xuất context & sinh câu trả lời.", "size": 16, "color": TEXT_WHITE, "level": 1},
        {"text": "• Model sử dụng: gpt-4o-mini (cho API chat server-side để tối ưu latency và chi phí).", "size": 16, "color": TEXT_WHITE, "level": 1},
        {"text": "• Cost-of-error: AI chỉ đề xuất, hiển thị nguồn kèm chip chỉ dẫn để người dùng kiểm chứng (Augmented AI).", "size": 16, "color": TEXT_WHITE, "level": 1},
        
        {"text": "Kịch bản Demo Live (1 case chuẩn + 1 case bẫy rủi ro):", "size": 18, "color": ACCENT_BLUE, "bold": True},
        {"text": "• Case 1 (Happy Path): Hỏi tóm tắt buổi hôm nay -> Bot nhận diện, tự nới scope ra cả buổi học, trả về kết quả tóm tắt mạch lạc, phân bổ số trang đối chiếu chính xác, click nhảy trang hoạt động mượt.", "size": 15, "color": TEXT_GRAY, "level": 1},
        {"text": "• Case 2 (Graceful Deflection): Hỏi ví dụ dự án AI ở slide 4 (thực tế slide 4 không có) -> Bot không bịa đặt ví dụ, chỉ tóm tắt 3 ý chính trên slide và hướng dẫn người học nới scope lên Cả môn để tìm ví dụ cụ thể.", "size": 15, "color": TEXT_GRAY, "level": 1}
    ]
    add_bullet_points(slide3, Inches(0.75), Inches(1.3), Inches(11.833), Inches(5.5), points3)

    # -------------------------------------------------------------------------
    # SLIDE 4: Kết quả đo
    # -------------------------------------------------------------------------
    slide4 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide4, BG_COLOR)
    create_title(slide4, "4. KẾT QUẢ ĐO VÀ QUALITY BAR ĐÃ CAM KẾT")
    
    points4 = [
        {"text": "Quality Bar Cam Kết (Chốt từ 23:59 N1):", "size": 18, "color": ACCENT_BLUE, "bold": True},
        {"text": "• NLU Intent Accuracy: >= 80% (đo trên 17 case) | OOD False Positive: 0/3 case.", "size": 16, "color": TEXT_WHITE, "level": 1},
        {"text": "• Answer Quality: Accuracy >= 1.5, Relevance >= 1.5, Completeness >= 1.4, Hallucination Rate <= 10%.", "size": 16, "color": TEXT_WHITE, "level": 1},
        
        {"text": "Kết quả Đo Lần Cuối (Tổng 38 case):", "size": 18, "color": ACCENT_BLUE, "bold": True},
        {"text": "• Bộ NLU Intent Classification: Đạt 14/17 câu (82.35%) -> ĐẠT QUALITY BAR! (OOD FPR: 0/3)", "size": 16, "color": ACCENT_GREEN, "bold": True, "level": 1},
        {"text": "• Bộ Answer Quality (GPT-4o-mini): Đạt 20/21 câu đạt -> ĐẠT TRỌN BỘ QUALITY BAR!", "size": 16, "color": ACCENT_GREEN, "bold": True, "level": 1},
        {"text": "- Accuracy trung bình: 1.71 (so với bar 1.5)", "size": 15, "color": TEXT_GRAY, "level": 2},
        {"text": "- Relevance trung bình: 1.90 (so với bar 1.5)", "size": 15, "color": TEXT_GRAY, "level": 2},
        {"text": "- Completeness trung bình: 1.62 (so với bar 1.4)", "size": 15, "color": TEXT_GRAY, "level": 2},
        {"text": "- Hallucination rate: 0.0% (so với bar <= 10%)", "size": 15, "color": TEXT_GRAY, "level": 2},
        
        {"text": "Phân tích nguyên nhân Ca thất bại lớn nhất (3/17 case NLU hỏng):", "size": 18, "color": ACCENT_BLUE, "bold": True},
        {"text": "• Các case lỗi: tc_011 ('slie 24'), tc_012 ('slde 26'), tc_016 ('sờ lai này') do người dùng gõ sai chính tả / từ lóng.", "size": 15, "color": TEXT_WHITE, "level": 1},
        {"text": "• Nguyên nhân: Regex khớp từ khóa quá cứng nhắc, dễ vỡ trước ngôn ngữ khẩu ngữ tự nhiên.", "size": 15, "color": TEXT_GRAY, "level": 2}
    ]
    add_bullet_points(slide4, Inches(0.75), Inches(1.3), Inches(11.833), Inches(5.5), points4)

    # -------------------------------------------------------------------------
    # SLIDE 5: User thật nói gì
    # -------------------------------------------------------------------------
    slide5 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide5, BG_COLOR)
    create_title(slide5, "5. PHẢN HỒI CỦA USER THẬT (VALIDATION LOG)")
    
    points5 = [
        {"text": "Nhật ký kiểm thử thực tế từ 5 người thử ngoài nhóm (Zone E403):", "size": 18, "color": ACCENT_BLUE, "bold": True},
        
        {"text": "1. Nguyễn Văn Hưng (Học viên ngoài nhóm):", "size": 16, "color": TEXT_WHITE, "bold": True, "level": 1},
        {"text": "• Quote nguyên văn: “bấm tóm tắt lần đầu thì k đc, nma bấm 'tạo lại câu trả lời' thì lại đc”", "size": 15, "color": ACCENT_ORANGE, "level": 2},
        {"text": "-> Thay đổi kỹ thuật: Tăng thời hạn timeout và thêm cơ chế tự động retry 2 lần phía API Route.", "size": 15, "color": TEXT_GRAY, "level": 2},
        
        {"text": "2. Trương Thảo Nguyên (Học viên ngoài nhóm):", "size": 16, "color": TEXT_WHITE, "bold": True, "level": 1},
        {"text": "• Quote nguyên văn: “Ý là t muốn hỏi về tất cả các trang... Nhưng t ấn chọn cả buổi Thì kết quả vẫn là trang 35”", "size": 15, "color": ACCENT_ORANGE, "level": 2},
        {"text": "-> Thay đổi kỹ thuật: Sửa lỗi kẹt scope, đồng bộ biến phạm vi chọn thủ công từ client trực tiếp vào RAG query.", "size": 15, "color": TEXT_GRAY, "level": 2},
        
        {"text": "3. Nguyễn Đăng Hưng (Học viên ngoài nhóm):", "size": 16, "color": TEXT_WHITE, "bold": True, "level": 1},
        {"text": "• Quote nguyên văn: Thích thú và thả emoji cười đắc ý khi hỏi ví dụ slide 4 và nhận phản hồi khước từ khéo léo.", "size": 15, "color": ACCENT_GREEN, "level": 2},
        {"text": "-> Giữ nguyên thiết kế: Khẳng định thiết kế từ chối thông minh (Graceful Deflection) hoạt động rất tốt.", "size": 15, "color": TEXT_GRAY, "level": 2}
    ]
    add_bullet_points(slide5, Inches(0.75), Inches(1.3), Inches(11.833), Inches(5.5), points5)

    # -------------------------------------------------------------------------
    # SLIDE 6: Nếu có thêm 1 tuần
    # -------------------------------------------------------------------------
    slide6 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide6, BG_COLOR)
    create_title(slide6, "6. ROADMAP PHÁT TRIỂN & BÀI HỌC KINH NGHIỆM")
    
    points6 = [
        {"text": "Các nhiệm vụ ưu tiên nếu có thêm 1 tuần (Roadmap):", "size": 18, "color": ACCENT_BLUE, "bold": True},
        {"text": "1. Nâng cấp bộ đoán Intent sang LLM Classifier (có Structured Output JSON Schema)", "size": 16, "color": TEXT_WHITE, "bold": True, "level": 1},
        {"text": "• Mục tiêu: Xóa bỏ 100% lỗi đoán sai do người dùng gõ sai chính tả, tiếng lóng ('slde', 'slie', 'sờ lai').", "size": 15, "color": TEXT_GRAY, "level": 2},
        
        {"text": "2. Cải thiện thiết kế responsive giao diện bảng biểu cho Mobile (Feedback của Yến)", "size": 16, "color": TEXT_WHITE, "bold": True, "level": 1},
        {"text": "• Mục tiêu: Thêm thuộc tính overflow-x-auto cho các component Markdown trả về, giúp người dùng cuộn ngang bảng deadline trên điện thoại dễ dàng.", "size": 15, "color": TEXT_GRAY, "level": 2},
        
        {"text": "3. Tích hợp RAG nâng cao (Hybrid Search & Re-ranking)", "size": 16, "color": TEXT_WHITE, "bold": True, "level": 1},
        {"text": "• Mục tiêu: Tự động tổng hợp và sinh takeaway chuẩn xác hơn dựa trên việc gộp các slide có độ tương quan cao.", "size": 15, "color": TEXT_GRAY, "level": 2},
        
        {"text": "Bài học lớn nhất từ dự án (Key Lesson):", "size": 18, "color": ACCENT_BLUE, "bold": True},
        {"text": "“Kiểm thử sớm bằng dữ liệu thực và các lỗi của người dùng thật (typos, tiếng lóng, scope lock) là con đường ngắn nhất để đưa sản phẩm AI thoát khỏi phòng lab và đáp ứng đúng kỳ vọng của thực tế.”", "size": 16, "color": ACCENT_GREEN, "bold": True, "level": 1}
    ]
    add_bullet_points(slide6, Inches(0.75), Inches(1.3), Inches(11.833), Inches(5.5), points6)

    # Save presentation
    output_path = "demo-slides.pptx"
    prs.save(output_path)
    print(f"PowerPoint slides generated successfully at {output_path}!")

if __name__ == "__main__":
    main()
